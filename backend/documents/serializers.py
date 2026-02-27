"""
Document 序列化器
"""

from rest_framework import serializers
from .models import Document
from shares.models import Share


class DocumentSerializer(serializers.ModelSerializer):
    """文档序列化器"""

    creator_name = serializers.CharField(source='creator.username', read_only=True)
    permission = serializers.SerializerMethodField()
    is_shared = serializers.SerializerMethodField()
    share_count = serializers.SerializerMethodField()

    class Meta:
        model = Document
        fields = [
            'id', 'title', 'type', 'file', 'creator', 'creator_name',
            'created_at', 'updated_at', 'version', 'permission',
            'is_shared', 'share_count',
            'is_deleted'
        ]
        read_only_fields = ['creator', 'version', 'created_at', 'updated_at']

    def get_permission(self, obj):
        """获取当前用户的权限"""
        request = self.context.get('request')
        if not request:
            return None

        # 如果是所有者
        if request.user == obj.creator:
            return 'owner'

        annotated_permission = getattr(obj, 'current_user_share_permission', None)
        if annotated_permission:
            return annotated_permission

        # 查询分享权限
        share = Share.objects.filter(
            document=obj,
            sharee=request.user,
            is_active=True
        ).first()

        if share:
            return share.permission

        return None

    def create(self, validated_data):
        """创建文档时自动设置创建者和空白文件"""
        request = self.context.get('request')
        if request and hasattr(request, 'user'):
            validated_data['creator'] = request.user

        # 先创建文档记录
        document = super().create(validated_data)

        # 根据文档类型创建空白文件
        doc_type = validated_data.get('type', 'word')
        file_ext = 'docx'
        content_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'

        if doc_type == 'cell':
            file_ext = 'xlsx'
            content_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        elif doc_type == 'slide':
            file_ext = 'pptx'
            content_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'

        # 创建空白文件
        from django.core.files.base import ContentFile
        import io

        if doc_type == 'word':
            from docx import Document as DocxDocument
            doc = DocxDocument()
            doc.add_paragraph('')
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            filename = f"{document.title}.{file_ext}"
            document.file.save(filename, ContentFile(buffer.read()), save=True)
        elif doc_type == 'cell':
            from openpyxl import Workbook
            wb = Workbook()
            ws = wb.active
            ws.title = "Sheet1"
            buffer = io.BytesIO()
            wb.save(buffer)
            buffer.seek(0)
            filename = f"{document.title}.{file_ext}"
            document.file.save(filename, ContentFile(buffer.read()), save=True)
        elif doc_type == 'slide':
            from pptx import Presentation
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[0])
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            filename = f"{document.title}.{file_ext}"
            document.file.save(filename, ContentFile(buffer.read()), save=True)

        return document

    def get_is_shared(self, obj):
        """检查文档是否被分享"""
        if hasattr(obj, 'is_shared_annotated'):
            return obj.is_shared_annotated
        return obj.shares.filter(is_active=True).exists()

    def get_share_count(self, obj):
        """获取文档被分享的次数"""
        if hasattr(obj, 'share_count_annotated'):
            return obj.share_count_annotated
        return obj.shares.filter(is_active=True).count()


class DocumentListSerializer(serializers.ModelSerializer):
    """文档列表序列化器（精简版）"""

    creator_name = serializers.CharField(source='creator.username', read_only=True)
    permission = serializers.SerializerMethodField()
    is_shared = serializers.SerializerMethodField()
    share_count = serializers.SerializerMethodField()

    class Meta:
        model = Document
        fields = [
            'id', 'title', 'type', 'creator_name', 'updated_at',
            'version', 'is_deleted', 'permission',
            'is_shared', 'share_count'
        ]

    def get_permission(self, obj):
        """获取当前用户的权限"""
        request = self.context.get('request')
        if not request:
            return None

        # 如果是所有者
        if request.user == obj.creator:
            return 'owner'

        annotated_permission = getattr(obj, 'current_user_share_permission', None)
        if annotated_permission:
            return annotated_permission

        # 查询分享权限
        share = Share.objects.filter(
            document=obj,
            sharee=request.user,
            is_active=True
        ).first()

        if share:
            return share.permission

        return None

    def get_is_shared(self, obj):
        """检查文档是否被分享"""
        if hasattr(obj, 'is_shared_annotated'):
            return obj.is_shared_annotated
        return obj.shares.filter(is_active=True).exists()

    def get_share_count(self, obj):
        """获取文档被分享的次数"""
        if hasattr(obj, 'share_count_annotated'):
            return obj.share_count_annotated
        return obj.shares.filter(is_active=True).count()
